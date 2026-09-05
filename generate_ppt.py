import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BLACK = RGBColor(0x01, 0x17, 0x11)
    ONYX = RGBColor(0x04, 0x25, 0x1c)
    EMERALD_ACCENT = RGBColor(0x34, 0xd3, 0x99)
    MINT_GREEN = RGBColor(0x10, 0xb9, 0x81)
    GOLD_TEXT = RGBColor(0xfb, 0xbf, 0x24)
    WHITE = RGBColor(0xff, 0xff, 0xff)
    GRAY_TEXT = RGBColor(0x9c, 0xa3, 0xaf)
    ROSE_TEXT = RGBColor(0xf8, 0x71, 0x71)

    slides_data = [
        # Slide 1
        {
            "title": "ANTIGRAVITY AG 3D",
            "subtitle": "AI Project Idea Generator & Capstone Architect for Engineering Students",
            "bullets": [
                "20-Slide Comprehensive Architectural & Mentoring Presentation Deck",
                "Profile Scoped: Beginner Student • Intern / Junior Developer • Final-Year Capstone",
                "Featuring AI Top Recommendation Engine, Data Platform Badges & 60s Demo Scripts",
                "Live Web Application: https://3hr-promptwars.vercel.app"
            ],
            "quote": "Engineering is the closest thing to magic that exists in the real world. — Elon Musk (CEO, Tesla & xAI)"
        },
        # Slide 2
        {
            "title": "The Student Capstone Dilemma",
            "subtitle": "Why Final-Year Engineering Students Struggle with Project Selection",
            "bullets": [
                "Over-ambitious Scope: Selecting enterprise-scale projects that stall halfway through the semester.",
                "Unreachable Data Dependencies: Relying on proprietary datasets or APIs behind expensive paywalls.",
                "Vague MVP Definitions: Lacking an explicit feature cut and weekly sequential build order.",
                "Judge Presentation Gaps: Inability to structure a crisp 60-second live demo script highlighting core AI."
            ],
            "quote": "AI will be the ultimate tool to expand human capability and accelerate scientific discovery. — Demis Hassabis (CEO, Google DeepMind)"
        },
        # Slide 3
        {
            "title": "The Antigravity AG Solution",
            "subtitle": "AI-Driven Profile Scoper & Feasibility Check Engine",
            "bullets": [
                "Input Parameters: Known skills, multi-domain preferences, available weeks, and hardware constraints.",
                "Smart Profile Adapter: Adjusts complexity dynamically for Beginners, Interns, and Final-Year Students.",
                "Feasibility First: Enforces 9 strict constraints ensuring 100% buildable projects without data paywalls.",
                "Instant Output: Generates 3 tailored project ideas with an explicit AI Top Recommendation Choice."
            ],
            "quote": "The web is the most accessible runtime on earth — shipping software in seconds changes everything. — Guillermo Rauch (CEO, Vercel)"
        },
        # Slide 4
        {
            "title": "Integrated Tech & AI Ecosystem Logos",
            "subtitle": "Powered by Industry-Standard Artificial Intelligence & Cloud Platforms",
            "bullets": [
                "OpenAI / ChatGPT API: Structured JSON prompt execution & natural language budget guidance.",
                "Google Gemini 2.5 Flash: High-reasoning capstone mentoring & architectural blueprint elaboration.",
                "Vercel Cloud Serverless: Instant global edge deployment with Python serverless functions.",
                "PyTorch & Scikit-Learn: Computer vision transfer learning & machine learning classifiers.",
                "FastAPI & React: Async REST backend architecture and modern 3D glassmorphic user interfaces.",
                "GitHub REST API: Public repository diff ingestion and static code analysis integration."
            ],
            "quote": "Building reliable, safe AI systems requires deep engineering discipline from day one. — Dario Amodei (CEO, Anthropic)"
        },
        # Slide 5
        {
            "title": "Target Profile Scoper Matrix",
            "subtitle": "Tailored Complexity Scoping across Student Experience Tiers",
            "bullets": [
                "Beginner Student: Pretrained models (MobileNet, OpenCV), 1-2 core features, straightforward scope.",
                "Intern / Junior Developer: Clean REST API endpoints, full-stack integration, industry standard tools.",
                "Final-Year Student: Multi-component ML pipelines, dataset validation, research novelty, judge ready.",
                "Dynamic Milestone Adjustment: Week-by-week build paths explicitly cite student's primary language."
            ],
            "quote": "Software is a great combination between artistry and engineering. — Bill Gates (Co-Founder, Microsoft)"
        },
        # Slide 6
        {
            "title": "Structured System Prompt Architecture",
            "subtitle": "The Mentoring Brain Behind Antigravity AG",
            "bullets": [
                "System Persona: Direct, practical senior engineering mentor — zero buzzwords or hype.",
                "Constraint Enforcer: Proposes synthetic fallback datasets if ideal data sources require paywalls.",
                "Ethical & Privacy Shield: Explicitly flags biometric, medical, and personal data privacy risks.",
                "Learning Gap Minimizer: Recommends smallest learning curve path matching student skills to domain."
            ],
            "quote": "Most of human knowledge is learned by doing and building hands-on systems. — Yann LeCun (Chief AI Scientist, Meta)"
        },
        # Slide 7
        {
            "title": "The 9-Field Project Output Schema",
            "subtitle": "Standardized Blueprint Structure for Every Generated Idea",
            "bullets": [
                "1. Title (Short, specific) | 2. One-Line Pitch | 3. Problem & Beneficiaries",
                "4. Core Features (3-5 MVP-first) | 5. Tech Stack with Justifications",
                "6. Data Sources & Data Platform Source (Kaggle, GitHub, WHO, OpenFDA)",
                "7. Milestones (Week-by-week path) | 8. Risk/Difficulty Flags & Ethical Warnings",
                "9. Stretch Feature for Bonus Marks | AI Recommendation Reason & Top Match Flag"
            ],
            "quote": "The biggest risk is not taking any risk in a world that is changing quickly. — Mark Zuckerberg (CEO, Meta)"
        },
        # Slide 8
        {
            "title": "AI Top Recommendation Engine",
            "subtitle": "Data-Driven Project Matching & Feasibility Scoring",
            "bullets": [
                "Match Analysis: Evaluates skill alignment, time budget feasibility, and presentation impact.",
                "Visual Highlight: Displays a gold crown banner (⭐ AI Top Recommended Choice) on the #1 match.",
                "Justification Statement: Provides explicit rationale of why AI recommends this specific project.",
                "Confidence Metrics: Verifies 100% dataset availability before presentation to judges."
            ],
            "quote": "Execution is everything. Ideas are cheap, building real products is hard. — Tim Cook (CEO, Apple)"
        },
        # Slide 9
        {
            "title": "Domain Blueprint 1 — Software Engineering",
            "subtitle": "DevMentor AI: Automated Code Reviewer for Student Repositories",
            "bullets": [
                "Pitch: An AI code-review assistant flagging code smells and suggesting refactors for student repos.",
                "Core Features: GitHub diff ingestion, static AST analysis, severity score dashboard, inline suggested diffs.",
                "Tech Stack: Python (FastAPI), React / Tailwind, GitHub REST API, Gemini LLM API.",
                "Data Platform: GitHub REST API & PyLint/ESLint open-source rule engines."
            ],
            "quote": "Simplicity is prerequisite for reliability. — Edsger W. Dijkstra"
        },
        # Slide 10
        {
            "title": "Domain Blueprint 2 — Data Science & ML",
            "subtitle": "SkillGap AI: Placement Predictor & Job-Readiness Analytics",
            "bullets": [
                "Pitch: ML analytics dashboard predicting student job readiness & recommending skill closing roadmaps.",
                "Core Features: Student profile feature engineering, XGBoost placement model, skill gap radar chart.",
                "Tech Stack: Python, scikit-learn, XGBoost, Pandas, Plotly, Streamlit / React.",
                "Data Platform: Kaggle Open Campus Placement Dataset & industry skill benchmarks."
            ],
            "quote": "Without data, you're just another person with an opinion. — W. Edwards Deming"
        },
        # Slide 11
        {
            "title": "Domain Blueprint 3 — NLP & Language Models",
            "subtitle": "Antigravity AG: Capstone Project Mentor & Architecture Advisor",
            "bullets": [
                "Pitch: NLP mentoring system matching student profiles to feasibility-checked project blueprints.",
                "Core Features: Skill profile NER extraction, embedding similarity search, 9-field project brief generator.",
                "Tech Stack: Python, FastAPI, Gemini API, Sentence-Transformers, Tailwind CSS.",
                "Data Platform: Curated Open Capstone Knowledge Base across 8 core CS specializations."
            ],
            "quote": "Language is the operating system of human thought. — Sam Altman"
        },
        # Slide 12
        {
            "title": "Domain Blueprint 4 — Computer Vision",
            "subtitle": "VisionTrack: Privacy-Preserving Classroom Attendance Logger",
            "bullets": [
                "Pitch: Campus face detection system logging attendance without saving raw facial embeddings.",
                "Core Features: MediaPipe face detection, head pose orientation score, CSV attendance sheet exporter.",
                "Tech Stack: Python, OpenCV, MediaPipe CPU engine, FastAPI, SQLite.",
                "Data Platform: MediaPipe Vision Engine & DAiSEE Classroom Video Dataset."
            ],
            "quote": "Vision is the art of seeing what is invisible to others. — Jonathan Swift"
        },
        # Slide 13
        {
            "title": "Domain Blueprint 5 — Healthcare & Medical ML",
            "subtitle": "ArogyaAssist: Multilingual Rural Symptom Checker & Triage Bot",
            "bullets": [
                "Pitch: Offline-first symptom guidance app classifying triage urgency for rural clinics with medical disclaimers.",
                "Core Features: Conversational symptom intake, probabilistic disease classifier, urgency triage badges.",
                "Tech Stack: Python, scikit-learn Decision Trees, FastAPI, Tailwind CSS.",
                "Data Platform: Kaggle Disease-Symptom Open Dataset & WHO Guidelines."
            ],
            "quote": "Healthcare is where technology meets human compassion directly. — Sundar Pichai"
        },
        # Slide 14
        {
            "title": "Domain Blueprint 6 — Fintech & Cyber Security",
            "subtitle": "SpendGuard Anomaly Detector & PhishShield Security Scanner",
            "bullets": [
                "SpendGuard: Isolation Forest financial transaction anomaly detector with AI budgeting nudges.",
                "PhishShield: Real-time phishing email header scanner, WHOIS domain age checker, & link expander.",
                "Data Platform Sources: Kaggle Synthetic Transaction Dataset & PhishTank Open API.",
                "Privacy First: All bank statements & email credentials processed locally in memory without disk storage."
            ],
            "quote": "Security is not a product, but a continuous process. — Bruce Schneier"
        },
        # Slide 15
        {
            "title": "Domain Blueprint 7 — Sustainability & IoT",
            "subtitle": "EcoSort: Edge Waste-Sorting Assistant & Recycling Classifier",
            "bullets": [
                "Pitch: Camera waste classifier guiding recyclable vs compostable disposal with gamified eco-scores.",
                "Core Features: MobileNet image classification, disposal bin recommendation, campus eco-leaderboard.",
                "Tech Stack: Python, PyTorch / Transfer Learning, FastAPI, Tailwind CSS.",
                "Data Platform: Kaggle TrashNet Open Dataset (2,500+ labeled images)."
            ],
            "quote": "Sustainability is no longer about doing less harm. It is about doing more good. — Jochen Zeitz"
        },
        # Slide 16
        {
            "title": "Section 4 Deep-Dive Mentor Blueprint",
            "subtitle": "Actionable Implementation & Sequential Build Order",
            "bullets": [
                "Minimal File Tree: Standardized clean project directory layout (backend/, frontend/, requirements.txt).",
                "MVP Feature Cut: Explicit list of features to build first vs features to explicitly skip.",
                "Domain-Specific Skips: Software Eng skips OAuth orgs; Vision skips 3D matching; Healthcare skips EHR.",
                "Sequential Build Order: Numbered step-by-step coding sequence preventing technical dependency blocks."
            ],
            "quote": "First, solve the problem. Then, write the code. — John Johnson"
        },
        # Slide 17
        {
            "title": "Technical Blockers Registry & 1-Line Fixes",
            "subtitle": "Pre-Emptive Debugging Solutions for Common Failure Points",
            "bullets": [
                "Blocker 1: Gemini API returning markdown code blocks -> Fix: Strip regex markdown fences before json.loads().",
                "Blocker 2: GitHub API rate limits -> Fix: Pass GitHub Personal Access Token in request headers.",
                "Blocker 3: CORS issues on serverless -> Fix: Add FastAPI CORSMiddleware with allow_origins=['*'].",
                "Blocker 4: Isolation Forest flagging fixed rent payments -> Fix: Add rule-based threshold overrides."
            ],
            "quote": "Debugging is twice as hard as writing the code in the first place. — Brian W. Kernighan"
        },
        # Slide 18
        {
            "title": "Data Platform Provenance & Privacy Ethics",
            "subtitle": "Transparent Data Grounding & Mandatory Disclaimers",
            "bullets": [
                "Data Platform Provenance: Every project explicitly badges its verified open data source.",
                "Privacy-First Processing: Camera feeds & financial CSVs processed strictly in RAM without logging.",
                "Medical Triage Safety: Prominent un-dismissable disclaimers: 'Educational Triage Tool Only'.",
                "Demographic Bias Removal: Protected attributes (gender, origin) explicitly dropped during ML training."
            ],
            "quote": "Ethics and privacy must be engineered into software from the very first line of code."
        },
        # Slide 19
        {
            "title": "Hackathon 60-Second Demo & 20-Min Wow Factor",
            "subtitle": "Winning Presentation Techniques for Evaluation Panels",
            "bullets": [
                "60-Second Judge Demo Script: Concise script highlighting input, processing pipeline, and results in < 60s.",
                "20-Minute Wow Factors: One-click PDF brief export, language switcher, or README.md generator.",
                "Judge Questions Ready: Prepared answers for model accuracy, dataset fallbacks, and deployment scaling."
            ],
            "quote": "A demo is worth a thousand slides. Clean execution wins hackathons."
        },
        # Slide 20
        {
            "title": "Live Deployment & Repository Access",
            "subtitle": "Antigravity AG 3D Capstone AI Mentor is Live in Production",
            "bullets": [
                "Live Production Web App: https://3hr-promptwars.vercel.app",
                "GitHub Repository: https://github.com/ishansharma16012007-bot/3hr-promptwars",
                "Localhost Engine: http://127.0.0.1:8000",
                "Built with FastAPI, Gemini 2.5 Flash, Tailwind CSS, 3D Glassmorphism, & Vercel Serverless."
            ],
            "quote": "Thank you! Built for engineering students pushing the boundaries of software innovation."
        }
    ]

    for slide_idx, data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout

        # Dark onyx background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLACK
        bg.line.fill.background()

        # Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = EMERALD_ACCENT

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
        tf_sub = sub_box.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = data["subtitle"]
        p_sub.font.size = Pt(16)
        p_sub.font.bold = True
        p_sub.font.color.rgb = GOLD_TEXT

        # Glassmorphic Content Card Background Shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.7), Inches(4.3))
        card.fill.solid()
        card.fill.fore_color.rgb = ONYX
        card.line.color.rgb = MINT_GREEN

        # Bullets Frame
        content_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(3.9))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True

        for b_idx, bullet in enumerate(data["bullets"]):
            p_b = tf_content.add_paragraph() if b_idx > 0 else tf_content.paragraphs[0]
            p_b.text = f"•  {bullet}"
            p_b.font.size = Pt(15)
            p_b.font.color.rgb = WHITE
            p_b.space_after = Pt(12)

        # Footer Quote Box
        if data.get("quote"):
            quote_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7))
            tf_q = quote_box.text_frame
            tf_q.word_wrap = True
            p_q = tf_q.paragraphs[0]
            p_q.text = f"💬  \"{data['quote']}\""
            p_q.font.size = Pt(13)
            p_q.font.italic = True
            p_q.font.color.rgb = MINT_GREEN

    output_path = "antigravity_presentation.pptx"
    prs.save(output_path)
    print(f"PowerPoint Presentation created successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
